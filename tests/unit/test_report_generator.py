"""Tests for the multi-format report generator and report router.

Covers:
- PDF output is non-empty and contains every required report section.
- DOCX output is a real OOXML document (python-docx can open it, contains
  headings, tables, and embedded figure images).
- TXT output is plain text with all required sections.
- PPTX output is a valid presentation.
- Defensive handling of missing / differently-keyed payloads (the key
  mismatches identified in the audit).
"""

import asyncio
import io
import math
import zipfile

import pytest

from backend.domain.entities.experiment import Experiment, ExperimentMetadata
from backend.services.report_generator import (
    ReportGenerator,
    generate_report_bytes,
)


def _make_pattern(start=10.0, stop=90.0, step=0.02):
    two_theta = [start + i * step for i in range(int((stop - start) / step))]
    peaks = [28.44, 47.30, 56.12, 69.13]
    intensity = []
    for t in two_theta:
        bg = 50 + 0.1 * t
        signal = sum(500 * math.exp(-0.5 * ((t - pk) / 0.15) ** 2) for pk in peaks)
        noise = (hash(str(t)) % 10 - 5) * 0.5
        intensity.append(bg + signal + noise)
    return two_theta, intensity


def _project_data():
    return {
        "name": "Si Reference Study",
        "material": "Si",
        "created_at": "2026-01-15",
        "status": "Complete",
    }


def _experiment_data():
    tt, intensity = _make_pattern()
    return {
        "name": "Si standard",
        "two_theta": tt,
        "intensity": intensity,
        "processed_pattern": {
            "two_theta": tt,
            "intensity": [i * 0.9 for i in intensity],
        },
        "detected_peaks": [
            {"two_theta": 28.44, "intensity": 500.0, "d_spacing": 3.135, "fwhm": 0.18},
            {"two_theta": 47.30, "intensity": 260.0, "d_spacing": 1.920, "fwhm": 0.20},
            {"two_theta": 56.12, "intensity": 180.0, "d_spacing": 1.638, "fwhm": 0.21},
        ],
        "candidate_phases": [
            {
                "rank": 1,
                "material_name": "Silicon",
                "material_formula": "Si",
                "match_score": 0.93,
                "fom": 0.12,
                "confidence": "High",
                "matched_peaks": 4,
                "total_reference_peaks": 5,
            },
            {
                "rank": 2,
                "material_name": "Quartz low",
                "material_formula": "SiO2",
                "match_score": 0.51,
                "fom": 2.40,
                "confidence": "Low",
                "matched_peaks": 2,
                "total_reference_peaks": 8,
            },
        ],
        "rietveld_results": {
            "status": "completed",
            "r_wp": 8.2,
            "r_p": 6.1,
            "r_exp": 5.4,
            "chi_squared": 2.30,
            "gof": 1.52,
            "iterations": 120,
            "parameters": {
                "scale": 1.02,
                "zero_shift": 0.001,
                "background_coeffs": [10.0, 1.0, 0.0, 0.0],
                "U": 0.006,
                "V": -0.002,
                "W": 0.008,
                "phase_fractions": [1.0],
            },
            "phases_used": [
                {
                    "name": "Silicon",
                    "formula": "Si",
                    "space_group": "Fd-3m",
                    "fraction": 1.0,
                    "lattice_params": {"a": 5.4310},
                    "lattice_param_uncertainties": {"a": 0.0002},
                }
            ],
            "patterns": {
                "two_theta": tt,
                "observed": intensity,
                "calculated": [i * 0.99 for i in intensity],
                "difference": [i * 0.01 for i in intensity],
            },
        },
        "pipeline_stages": [
            {"name": "background_correction", "status": "completed", "duration_seconds": 1.2},
            {"name": "peak_detection", "status": "completed", "duration_seconds": 0.4},
            {"name": "rietveld_refinement", "status": "completed", "duration_seconds": 3.0},
        ],
        "wavelength": 1.5406,
        "metadata": {
            "instrument": "Bruker D8 Advance",
            "radiation_type": "Cu",
            "temperature_k": 298.0,
            "scan_range_2theta": [10.0, 90.0],
            "step_size_2theta": 0.02,
            "scan_time_seconds": 300.0,
            "notes": "Synthesized by solid-state reaction.",
        },
    }


class TestPdfReport:
    def test_pdf_nonempty_and_valid_header(self):
        pdf = ReportGenerator().generate_report_bytes(_project_data(), _experiment_data())
        assert len(pdf) > 1000
        assert pdf[:4] == b"%PDF"

    def test_pdf_contains_all_required_sections(self):
        pdf = ReportGenerator().generate_report_bytes(_project_data(), _experiment_data())
        required = [
            "Project Information",
            "Sample Information",
            "Experimental Conditions",
            "Data Summary",
            "Detected Peaks",
            "Processing Workflow",
            "Phase Identification Results",
            "Rietveld Refinement Summary",
            "Refinement Statistics",
            "Phase Fractions",
            "Refined Parameters",
            "Scientific Conclusions",
            "Methodology",
            "References",
            "Figures",
        ]
        for section in required:
            assert section.encode() in pdf, f"missing section in PDF: {section}"

    def test_pdf_contains_peak_and_phase_values(self):
        pdf = ReportGenerator().generate_report_bytes(_project_data(), _experiment_data())
        assert b"28.440" in pdf
        assert b"3.1350" in pdf
        assert b"Silicon" in pdf
        assert b"1.5406" in pdf

    def test_pdf_generate_to_disk(self, tmp_path):
        out = tmp_path / "report.pdf"
        path = ReportGenerator().generate_report(_project_data(), _experiment_data(), str(out))
        assert out.exists()
        assert out.read_bytes()[:4] == b"%PDF"
        assert str(out) == path


class TestDocxReport:
    def test_docx_is_real_ooxml_zip(self):
        from docx import Document

        docx_bytes = ReportGenerator().generate_docx_bytes(_project_data(), _experiment_data())
        assert docx_bytes[:4] == b"PK\x03\x04"
        assert zipfile.is_zipfile(io.BytesIO(docx_bytes))
        with zipfile.ZipFile(io.BytesIO(docx_bytes)) as zf:
            names = zf.namelist()
            assert "[Content_Types].xml" in names
            assert any(n.endswith("document.xml") for n in names)
        # python-docx must be able to open it.
        doc = Document(io.BytesIO(docx_bytes))
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        for expected in ("1. Project Information", "6. Phase Identification Results",
                         "10. Methodology", "11. References", "12. Figures"):
            assert expected in headings, f"missing DOCX heading: {expected}"
        assert len(doc.tables) >= 8

    def test_docx_contains_embedded_figure_images(self):
        from docx import Document

        docx_bytes = ReportGenerator().generate_docx_bytes(_project_data(), _experiment_data())
        doc = Document(io.BytesIO(docx_bytes))
        assert len(doc.inline_shapes) >= 2

    def test_docx_contains_peak_table_data(self):
        from docx import Document

        docx_bytes = ReportGenerator().generate_docx_bytes(_project_data(), _experiment_data())
        doc = Document(io.BytesIO(docx_bytes))
        all_text = []
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_text.append(cell.text)
        assert "Silicon" in all_text
        assert "28.440" in all_text
        assert any("1.52" in text for text in all_text)  # GoF


class TestTxtReport:
    def test_txt_is_utf8_plain_text_with_sections(self):
        txt = ReportGenerator().generate_txt_bytes(_project_data(), _experiment_data())
        text = txt.decode("utf-8")
        assert "MatPilot Scientific Analysis Report" in text
        for marker in ("PROJECT INFORMATION", "SAMPLE INFORMATION",
                       "EXPERIMENTAL CONDITIONS", "DATA SUMMARY",
                       "DETECTED PEAKS", "PHASE IDENTIFICATION RESULTS",
                       "RIETVELD REFINEMENT SUMMARY", "SCIENTIFIC CONCLUSIONS",
                       "METHODOLOGY", "REFERENCES"):
            assert marker in text, f"missing TXT section: {marker}"


class TestPptxReport:
    def test_pptx_is_valid_presentation(self):
        from pptx import Presentation

        pptx_bytes = ReportGenerator().generate_pptx_bytes(_project_data(), _experiment_data())
        assert pptx_bytes[:4] == b"PK\x03\x04"
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) >= 10

    def test_pptx_contains_figure_pictures(self):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        pptx_bytes = ReportGenerator().generate_pptx_bytes(_project_data(), _experiment_data())
        prs = Presentation(io.BytesIO(pptx_bytes))
        picture_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_count += 1
        assert picture_count >= 3


class TestDefensiveHandling:
    def test_minimal_payload_does_not_crash(self):
        gen = ReportGenerator()
        project = {"name": "Minimal"}
        for exp in (
            {},
            {"name": "Empty"},
            {"two_theta": [10, 20], "intensity": []},
        ):
            assert gen.generate_report_bytes(project, exp)[:4] == b"%PDF"
            assert gen.generate_docx_bytes(project, exp)[:4] == b"PK\x03\x04"
            assert "MatPilot" in gen.generate_txt_bytes(project, exp).decode("utf-8")
            assert gen.generate_pptx_bytes(project, exp)[:4] == b"PK\x03\x04"

    def test_mismatched_processed_pattern_keys_normalized(self):
        """The pipeline stores the processed pattern as a nested dict
        (experiment._processed_pattern); the generator previously looked for
        flat `processed_two_theta`/`processed_intensity` keys."""
        tt, intensity = _make_pattern()
        exp = {
            "two_theta": tt,
            "intensity": intensity,
            "processed_pattern": {"two_theta": tt, "intensity": [i * 0.8 for i in intensity]},
        }
        gen = ReportGenerator()
        canonical = gen._canonicalize_experiment(exp)
        assert len(canonical["processed_two_theta"]) == len(tt)
        assert len(canonical["processed_intensity"]) == len(tt)

    def test_rietveld_actual_pipeline_keys_used(self):
        """The pipeline stores Rietveld output under `parameters` / `patterns`
        / `phases_used`; the generator must derive its refined-parameter and
        phase-fraction tables from that shape."""
        exp = _experiment_data()
        gen = ReportGenerator()
        canonical = gen._canonicalize_experiment(exp)
        rt = canonical["rietveld_results"]
        assert len(rt["phases"]) == 1
        assert rt["phases_used_count"] == 1
        assert any("Scale factor" in p["name"] for p in rt["refined_parameters"])
        assert any("a \u2014 Silicon" in p["name"] for p in rt["refined_parameters"])
        assert rt["calculated_two_theta"]

    def test_rietveld_legacy_keys_still_work(self):
        tt, intensity = _make_pattern()
        exp = {
            "two_theta": tt,
            "intensity": intensity,
            "rietveld_results": {
                "r_wp": 7.0,
                "gof": 1.3,
                "calculated_two_theta": tt,
                "calculated_intensity": [i * 0.99 for i in intensity],
                "difference_intensity": [i * 0.01 for i in intensity],
                "phases": [{"name": "A", "formula": "AB", "fraction": 0.5}],
                "refined_parameters": [
                    {"name": "a", "initial": 1.0, "refined": 1.1, "uncertainty": 0.01}
                ],
            },
        }
        gen = ReportGenerator()
        pdf = gen.generate_report_bytes(_project_data(), exp)
        assert b"7.0000" in pdf
        assert b"Refined Parameters" in pdf


class TestDispatcher:
    def test_generate_report_bytes_dispatch(self):
        project = _project_data()
        exp = _experiment_data()
        assert generate_report_bytes("pdf", project, exp)[:4] == b"%PDF"
        assert generate_report_bytes("docx", project, exp)[:4] == b"PK\x03\x04"
        assert "MatPilot" in generate_report_bytes("txt", project, exp).decode("utf-8")
        assert generate_report_bytes("pptx", project, exp)[:4] == b"PK\x03\x04"

    def test_dispatcher_rejects_unknown_format(self):
        with pytest.raises(ValueError):
            generate_report_bytes("html", _project_data(), _experiment_data())


class TestReportRouter:
    def _seed_experiment(self, app):
        tt, intensity = _make_pattern()
        exp = Experiment(
            name="Si standard",
            wavelength_angstrom=1.5406,
            raw_two_theta=tt,
            raw_intensity=intensity,
            detected_peaks=[
                {"two_theta": 28.44, "intensity": 500.0, "d_spacing": 3.135, "fwhm": 0.18}
            ],
            candidate_phases=[
                {"rank": 1, "material_name": "Silicon", "material_formula": "Si",
                 "match_score": 0.93, "confidence": "High"}
            ],
            rietveld_results={
                "r_wp": 8.2, "r_p": 6.1, "r_exp": 5.4, "chi_squared": 2.3,
                "gof": 1.52, "iterations": 120,
                "parameters": {
                    "scale": 1.02, "zero_shift": 0.001, "background_coeffs": [10, 1, 0, 0],
                    "U": 0.006, "V": -0.002, "W": 0.008, "phase_fractions": [1.0],
                },
                "phases_used": [{"name": "Silicon", "formula": "Si", "fraction": 1.0}],
                "patterns": {
                    "two_theta": tt, "observed": intensity,
                    "calculated": [i * 0.99 for i in intensity],
                    "difference": [i * 0.01 for i in intensity],
                },
            },
            pipeline_stages=[
                {"name": "background_correction", "status": "completed", "duration_seconds": 1.0}
            ],
            metadata=ExperimentMetadata(
                instrument="Bruker D8 Advance",
                radiation_type="Cu",
                temperature_k=298.0,
                scan_range_2theta=[10.0, 90.0],
                step_size_2theta=0.02,
            ),
        )
        exp._processed_pattern = {"two_theta": tt, "intensity": intensity}
        asyncio.run(app.state.container.uow.experiments.add(exp))
        return str(exp.id)

    def test_router_returns_pdf(self, client, app):
        exp_id = self._seed_experiment(app)
        resp = client.post(f"/report/generate/{exp_id}")
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith("application/pdf")
        assert resp.content[:4] == b"%PDF"

    def test_router_returns_real_docx(self, client, app):
        from docx import Document

        exp_id = self._seed_experiment(app)
        resp = client.post(f"/report/generate/{exp_id}?format=docx")
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        doc = Document(io.BytesIO(resp.content))
        headings = [p.text for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert "10. Methodology" in headings

    def test_router_returns_txt(self, client, app):
        exp_id = self._seed_experiment(app)
        resp = client.post(f"/report/generate/{exp_id}?format=txt")
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith("text/plain")
        assert "METHODOLOGY" in resp.text

    def test_router_returns_pptx(self, client, app):
        from pptx import Presentation

        exp_id = self._seed_experiment(app)
        resp = client.post(f"/report/generate/{exp_id}?format=pptx")
        assert resp.status_code == 200
        assert resp.headers["content-type"].startswith(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        prs = Presentation(io.BytesIO(resp.content))
        assert len(prs.slides) >= 10

    def test_router_unsupported_format_returns_400(self, client, app):
        exp_id = self._seed_experiment(app)
        resp = client.post(f"/report/generate/{exp_id}?format=html")
        assert resp.status_code == 400

    def test_router_missing_experiment_returns_404(self, client):
        resp = client.post("/report/generate/00000000-0000-0000-0000-000000000000")
        assert resp.status_code == 404
