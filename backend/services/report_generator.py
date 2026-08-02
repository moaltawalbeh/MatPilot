"""Multi-format scientific report generator for MatPilot XRD analysis.

Generates publication-quality reports as:
- PDF  (ReportLab + matplotlib)
- DOCX (python-docx, real OOXML)
- TXT  (plain text)
- PPTX (python-pptx)

The generator is defensive: it normalizes the various experiment/payload key
shapes produced by the pipeline (processed patterns, Rietveld results,
instrument metadata, detected peaks, candidate phases) into a stable internal
structure so a missing or differently-keyed field never crashes generation.
"""

import io
import logging
import os
import warnings
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import (
    BaseDocTemplate,
    Frame,
    Image,
    NextPageTemplate,
    PageBreak,
    PageTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
)

# Suppress noisy matplotlib warnings (font discovery, missing glyphs, etc.).
warnings.filterwarnings("ignore", category=UserWarning, module="matplotlib")
warnings.filterwarnings("ignore", category=UserWarning, message=r"findfont.*")
warnings.filterwarnings("ignore", category=UserWarning, message=r"Glyph .* missing.*")
logging.getLogger("matplotlib").setLevel(logging.ERROR)

try:  # python-docx
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches as _DocxInches
    from docx.shared import Pt as _DocxPt

    _DOCX_AVAILABLE = True
except Exception:  # pragma: no cover
    _DOCX_AVAILABLE = False

try:  # python-pptx
    from pptx import Presentation
    from pptx.util import Inches as _PptxInches
    from pptx.util import Pt as _PptxPt

    _PPTX_AVAILABLE = True
except Exception:  # pragma: no cover
    _PPTX_AVAILABLE = False


ACCENT_COLOR = colors.HexColor("#f97316")
ACCENT_LIGHT = colors.HexColor("#fff7ed")
HEADER_BG = colors.HexColor("#1a1a1a")
TEXT_COLOR = colors.HexColor("#1f2937")
MUTED_COLOR = colors.HexColor("#6b7280")
TABLE_HEADER_BG = colors.HexColor("#f97316")
TABLE_ALT_ROW = colors.HexColor("#fff7ed")
TABLE_HEADER_TEXT = colors.white
BORDER_COLOR = colors.HexColor("#e5e7eb")

FIGURE_WIDTH_CM = 16
FIGURE_HEIGHT_CM = 9

PIPELINE_STAGE_LABELS = {
    "background_correction": "Background Correction",
    "ka2_stripping": "Kα₂ Stripping",
    "noise_reduction": "Noise Reduction",
    "intensity_normalization": "Intensity Normalization",
    "peak_detection": "Peak Detection",
    "phase_identification": "Phase Identification",
    "candidate_selection": "Candidate Selection",
    "rietveld_refinement": "Rietveld Refinement",
}

REFERENCES = [
    "Cullity, B. D.; Stock, S. R. Elements of X-Ray Diffraction, 3rd ed.; "
    "Prentice Hall: Upper Saddle River, NJ, 2001.",
    "Rietveld, H. M. A Profile Refinement Method for Nuclear and Magnetic "
    "Structures. J. Appl. Crystallogr. 1969, 2, 65-71.",
    "McCusker, L. B.; Von Dreele, R. B.; Cox, D. E.; Louer, D.; Scardi, P. "
    "Rietveld Refinement Guidelines. J. Appl. Crystallogr. 1999, 32, 36-50.",
    "International Centre for Diffraction Data. Powder Diffraction File; "
    "Newtown Square, PA.",
    "Smith, D. K.; Nichols, M. C.; Zolensky, M. E. POWD10: A FORTRAN IV "
    "Program for Calculating X-ray Powder Diffraction Patterns.",
    "Savitzky, A.; Golay, M. J. E. Smoothing and Differentiation of Data by "
    "Simplified Least Squares Procedures. Anal. Chem. 1964, 36, 1627-1639.",
    "Toby, B. H. EXPGUI, a Graphical User Interface for GSAS. "
    "J. Appl. Crystallogr. 2001, 34, 210-213.",
]

FORMAT_MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "txt": "text/plain; charset=utf-8",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _safe_float(value: Any, default: Any = "N/A") -> Any:
    if value is None:
        return default
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _fmt(value: Any, fmt_spec: str = ".4f", default: str = "N/A") -> str:
    v = _safe_float(value, None)
    if v is None:
        return default
    try:
        return f"{float(v):{fmt_spec}}"
    except (TypeError, ValueError):
        return default


def _fmt_pct(value: Any, default: str = "N/A") -> str:
    v = _safe_float(value, None)
    if v is None:
        return default
    try:
        return f"{float(v):.1f}%"
    except (TypeError, ValueError):
        return default


def _fmt_int(value: Any, default: str = "N/A") -> str:
    v = _safe_float(value, None)
    if v is None:
        return default
    try:
        return str(int(round(float(v))))
    except (TypeError, ValueError):
        return default


def _as_list(value: Any) -> List[Any]:
    if value is None:
        return []
    if isinstance(value, (list, tuple)):
        return list(value)
    if hasattr(value, "tolist"):
        return list(value.tolist())
    try:
        return list(value)
    except TypeError:
        return []


def _as_float_list(value: Any) -> List[float]:
    out = []
    for v in _as_list(value):
        f = _safe_float(v, None)
        if f is not None:
            out.append(float(f))
    return out


def _register_unicode_fonts() -> Dict[str, str]:
    """Register a Unicode TTF font with ReportLab so Greek/subscript glyphs render.

    Falls back to the standard Helvetica family when matplotlib fonts are
    unavailable (e.g. bare CI environments).
    """
    fonts = {"normal": "Helvetica", "bold": "Helvetica-Bold", "oblique": "Helvetica-Oblique"}
    try:
        import matplotlib.font_manager as fm

        normal = fm.findfont(fm.FontProperties(family="DejaVu Sans", weight="normal"))
        bold = fm.findfont(fm.FontProperties(family="DejaVu Sans", weight="bold"))
        oblique = fm.findfont(fm.FontProperties(family="DejaVu Sans", weight="oblique"))
        if normal and normal.lower().endswith(".ttf"):
            pdfmetrics.registerFont(TTFont("DejaVu", normal))
            fonts["normal"] = "DejaVu"
        if bold and bold.lower().endswith(".ttf"):
            pdfmetrics.registerFont(TTFont("DejaVu-Bold", bold))
            fonts["bold"] = "DejaVu-Bold"
        if oblique and oblique.lower().endswith(".ttf"):
            pdfmetrics.registerFont(TTFont("DejaVu-Oblique", oblique))
            fonts["oblique"] = "DejaVu-Oblique"
        if fonts["normal"] == "DejaVu":
            pdfmetrics.registerFontFamily(
                "DejaVu",
                normal=fonts["normal"],
                bold=fonts["bold"],
                italic=fonts["oblique"],
                boldItalic=fonts["bold"],
            )
    except Exception:
        pass
    return fonts


class _NumberedCanvas:
    """Helper to add page numbers to the canvas."""

    def __init__(self, doc: BaseDocTemplate, title: str, fonts: Dict[str, str]):
        self._doc = doc
        self._title = title
        self._fonts = fonts
        self._saved_page_states: List[Dict[str, Any]] = []

    def __call__(self, canvas, doc):
        canvas.saveState()
        canvas.setFont(self._fonts["normal"], 8)
        canvas.setFillColor(MUTED_COLOR)
        canvas.drawString(
            doc.leftMargin,
            1.2 * cm,
            f"{self._title}",
        )
        canvas.drawRightString(
            A4[0] - doc.rightMargin,
            1.2 * cm,
            f"Page {doc.page}",
        )
        canvas.setStrokeColor(BORDER_COLOR)
        canvas.setLineWidth(0.5)
        canvas.line(
            doc.leftMargin,
            1.6 * cm,
            A4[0] - doc.rightMargin,
            1.6 * cm,
        )
        canvas.restoreState()


class ReportGenerator:
    """Generates publication-quality reports for XRD analysis in multiple formats."""

    def __init__(self):
        self._fonts = _register_unicode_fonts()
        self._styles = getSampleStyleSheet()
        self._register_custom_styles()
        self._story: List[Any] = []
        self._toc_entries: List[Dict[str, Any]] = []
        self._figures: List[Dict[str, Any]] = []
        self._figure_counter = 0
        self._section_counter = 0
        self._page_width = A4[0]
        self._page_height = A4[1]
        self._margin_left = 2.2 * cm
        self._margin_right = 2.2 * cm
        self._margin_top = 2.5 * cm
        self._margin_bottom = 2.5 * cm

    # ------------------------------------------------------------------
    # Canonicalization (defensive normalization of experiment data)
    # ------------------------------------------------------------------

    def _canonicalize_experiment(self, experiment_data: Dict[str, Any]) -> Dict[str, Any]:
        """Normalize the experiment payload into a stable internal structure.

        Supports both the keys used by the pipeline / domain entities and the
        historically-different keys the generator itself used to read, so a
        missing or re-keyed field never crashes report generation.
        """
        exp = dict(experiment_data or {})

        # Wavelength: accept wavelength, wavelength_angstrom, metadata value.
        wl = exp.get("wavelength", exp.get("wavelength_angstrom"))
        if wl is None:
            wl = 1.5406
        wl = _safe_float(wl, 1.5406)
        exp["wavelength"] = float(wl)
        exp["wavelength_angstrom"] = float(wl)

        # Processed pattern: accept flat keys or a nested dict (e.g. the
        # pipeline's `experiment._processed_pattern = {"two_theta": ..., "intensity": ...}`).
        processed = exp.get("processed_pattern")
        if isinstance(processed, dict):
            if not exp.get("processed_two_theta"):
                exp["processed_two_theta"] = _as_list(processed.get("two_theta", []))
            if not exp.get("processed_intensity"):
                exp["processed_intensity"] = _as_list(processed.get("intensity", []))
        if not exp.get("processed_two_theta"):
            exp["processed_two_theta"] = []
        if not exp.get("processed_intensity"):
            exp["processed_intensity"] = []
        exp["processed_two_theta"] = _as_float_list(exp["processed_two_theta"])
        exp["processed_intensity"] = _as_float_list(exp["processed_intensity"])

        # Raw pattern data.
        exp["two_theta"] = _as_float_list(exp.get("two_theta"))
        exp["intensity"] = _as_float_list(exp.get("intensity"))

        # Instrument metadata.
        meta = exp.get("metadata")
        if isinstance(meta, dict):
            if not exp.get("instrument") and meta.get("instrument"):
                exp["instrument"] = str(meta.get("instrument"))
            if not exp.get("radiation_type") and meta.get("radiation_type"):
                exp["radiation_type"] = str(meta.get("radiation_type"))
            if exp.get("wavelength_angstrom") in (1.5406,) and meta.get("wavelength_angstrom"):
                exp["wavelength_angstrom"] = float(meta.get("wavelength_angstrom"))
                exp["wavelength"] = float(meta.get("wavelength_angstrom"))
            if not exp.get("scan_range") and meta.get("scan_range_2theta"):
                exp["scan_range"] = _as_float_list(meta.get("scan_range_2theta"))
            if not exp.get("step_size") and meta.get("step_size_2theta"):
                exp["step_size"] = _safe_float(meta.get("step_size_2theta"), None)
            if not exp.get("temperature_k") and meta.get("temperature_k"):
                exp["temperature_k"] = _safe_float(meta.get("temperature_k"), None)
            if not exp.get("scan_time_seconds") and meta.get("scan_time_seconds"):
                exp["scan_time_seconds"] = _safe_float(meta.get("scan_time_seconds"), None)
            if not exp.get("notes") and meta.get("notes"):
                exp["notes"] = str(meta.get("notes"))

        # Detected peaks: canonical keys + computed relative intensity / esd.
        exp["detected_peaks"] = self._canonical_peaks(exp)

        # Candidate phases.
        phases = exp.get("candidate_phases") or []
        exp["candidate_phases"] = [p for p in phases if isinstance(p, dict)]

        # Pipeline stages.
        stages = exp.get("pipeline_stages") or []
        exp["pipeline_stages"] = [s for s in stages if isinstance(s, dict)]

        # Rietveld results.
        exp["rietveld_results"] = self._canonical_rietveld(exp.get("rietveld_results"))

        # Radiation display name.
        exp["radiation_type_display"] = self._radiation_display(
            exp.get("radiation_type"), exp.get("wavelength")
        )

        return exp

    def _canonical_peaks(self, exp: Dict[str, Any]) -> List[Dict[str, Any]]:
        peaks = exp.get("detected_peaks") or []
        canonical = []
        for p in peaks:
            if not isinstance(p, dict):
                continue
            pos = p.get("two_theta", p.get("position", p.get("pos")))
            if pos is None:
                continue
            canonical.append({
                "two_theta": pos,
                "intensity": _safe_float(p.get("intensity", p.get("height")), None),
                "d_spacing": p.get("d_spacing", p.get("dspacing")),
                "fwhm": p.get("fwhm"),
                "area": p.get("area"),
                "esd": p.get("esd", p.get("uncertainty", p.get("position_error"))),
            })

        intensities = [
            float(c["intensity"]) for c in canonical
            if c["intensity"] is not None and float(c["intensity"]) > 0
        ]
        max_int = max(intensities) if intensities else 1.0
        if max_int <= 0:
            max_int = 1.0
        for c in canonical:
            if c["intensity"] is not None:
                c["relative_intensity"] = float(c["intensity"]) / max_int * 100.0
            else:
                c["relative_intensity"] = None
        return canonical

    def _canonical_rietveld(self, rietveld: Any) -> Optional[Dict[str, Any]]:
        if not isinstance(rietveld, dict):
            return None
        rt = dict(rietveld)

        params = rt.get("parameters")
        if not isinstance(params, dict):
            params = {}

        patterns = rt.get("patterns")
        if not isinstance(patterns, dict):
            patterns = {}

        # Normalize phase fractions / phases_used -> canonical `phases` list.
        phases = rt.get("phases")
        if not isinstance(phases, list):
            phases = []
        phases_used = rt.get("phases_used")
        phases_used_count = None
        if isinstance(phases_used, int):
            phases_used_count = phases_used
        elif isinstance(phases_used, list):
            phases_used_count = len(phases_used)
            if not phases:
                phases = [
                    {
                        "name": p.get("name") or p.get("formula") or f"Phase {i + 1}",
                        "formula": p.get("formula", ""),
                        "space_group": p.get("space_group", ""),
                        "fraction": p.get("fraction"),
                        "lattice_params": p.get("lattice_params", {}),
                        "lattice_param_uncertainties": p.get("lattice_param_uncertainties", {}),
                    }
                    for i, p in enumerate(phases_used)
                    if isinstance(p, dict)
                ]
        elif not phases:
            fractions = params.get("phase_fractions")
            if isinstance(fractions, list) and fractions:
                phases = [
                    {"name": f"Phase {i + 1}", "formula": "", "fraction": f}
                    for i, f in enumerate(fractions)
                ]
                phases_used_count = len(fractions)
        if phases_used_count is None and phases:
            phases_used_count = len(phases)
        rt["phases"] = phases
        rt["phases_used_count"] = phases_used_count

        # Build refined_parameters from the flat `parameters` dict + uncertainties.
        refined_params = rt.get("refined_parameters")
        if not isinstance(refined_params, list):
            refined_params = []
        if not refined_params:
            unc_map = rt.get("parameter_uncertainties")
            if not isinstance(unc_map, dict):
                unc_map = {}

            def _add(name: str, value: Any, unc: Any = None) -> None:
                refined_params.append({
                    "name": name,
                    "initial": None,
                    "refined": value,
                    "uncertainty": unc,
                })

            _add("Scale factor", params.get("scale"), unc_map.get("scale"))
            _add("Zero shift (2\u03b8)", params.get("zero_shift"), unc_map.get("zero_shift"))
            bg = params.get("background_coeffs")
            if isinstance(bg, list):
                for i, b in enumerate(bg):
                    _add(f"Background coefficient {i + 1}", b)
            _add("Peak shape U", params.get("U"), unc_map.get("U"))
            _add("Peak shape V", params.get("V"), unc_map.get("V"))
            _add("Peak shape W", params.get("W"), unc_map.get("W"))
            fracs = params.get("phase_fractions")
            if isinstance(fracs, list):
                for i, f in enumerate(fracs):
                    label = "Phase fraction"
                    if i < len(phases) and phases[i].get("name"):
                        label = f"Fraction \u2014 {phases[i].get('name')}"
                    _add(label, f)
            if isinstance(phases_used, list):
                for pi, phase in enumerate(phases_used):
                    if not isinstance(phase, dict):
                        continue
                    lattice = phase.get("lattice_params") or {}
                    unc_lat = phase.get("lattice_param_uncertainties") or {}
                    pname = phase.get("name") or phase.get("formula") or f"Phase {pi + 1}"
                    for key in ("a", "b", "c", "alpha", "beta", "gamma", "volume"):
                        if key in lattice:
                            _add(f"{key} \u2014 {pname}", lattice[key], unc_lat.get(key))
        rt["refined_parameters"] = refined_params

        # Normalize fit patterns -> legacy keys the figure builder understands.
        if patterns:
            ctt = patterns.get("two_theta", patterns.get("calculated_two_theta", []))
            calc = patterns.get("calculated", patterns.get("calculated_intensity", []))
            diff = patterns.get("difference", patterns.get("difference_intensity", []))
            obs = patterns.get("observed", [])
            rt["calculated_two_theta"] = _as_float_list(ctt)
            rt["calculated_intensity"] = _as_float_list(calc)
            rt["difference_intensity"] = _as_float_list(diff)
            rt["observed_intensity"] = _as_float_list(obs)

        return rt

    def _radiation_display(self, radiation_type: Any, wavelength: Any) -> str:
        wl = _safe_float(wavelength, None)
        if radiation_type:
            rt = str(radiation_type)
            if rt in ("Cu", "Mo", "Co", "Fe", "Cr"):
                return f"{rt} K\u03b1"
            return rt
        if wl is not None:
            for wl_val, name in ((1.5406, "Cu K\u03b1"), (1.5419, "Cu K\u03b1 (avg)")):
                if abs(wl - wl_val) < 0.01:
                    return name
        return "X-ray"

    # ------------------------------------------------------------------
    # Styles
    # ------------------------------------------------------------------

    def _register_custom_styles(self):
        self._styles.add(ParagraphStyle(
            "CoverTitle",
            parent=self._styles["Title"],
            fontName=self._fonts["bold"],
            fontSize=28,
            leading=34,
            textColor=colors.white,
            alignment=TA_CENTER,
            spaceAfter=12,
        ))
        self._styles.add(ParagraphStyle(
            "CoverSubtitle",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=14,
            leading=18,
            textColor=colors.HexColor("#f97316"),
            alignment=TA_CENTER,
            spaceAfter=6,
        ))
        self._styles.add(ParagraphStyle(
            "CoverInfo",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=11,
            leading=15,
            textColor=colors.HexColor("#d1d5db"),
            alignment=TA_CENTER,
        ))
        self._styles.add(ParagraphStyle(
            "SectionHeading",
            parent=self._styles["Heading1"],
            fontName=self._fonts["bold"],
            fontSize=16,
            leading=20,
            textColor=colors.HexColor("#f97316"),
            spaceBefore=18,
            spaceAfter=10,
            borderWidth=0,
            borderColor=ACCENT_COLOR,
            borderPadding=0,
        ))
        self._styles.add(ParagraphStyle(
            "SubHeading",
            parent=self._styles["Heading2"],
            fontName=self._fonts["bold"],
            fontSize=12,
            leading=16,
            textColor=TEXT_COLOR,
            spaceBefore=10,
            spaceAfter=6,
        ))
        self._styles.add(ParagraphStyle(
            "BodyText2",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=10,
            leading=14,
            textColor=TEXT_COLOR,
            spaceAfter=6,
        ))
        self._styles.add(ParagraphStyle(
            "TableCaption",
            parent=self._styles["Normal"],
            fontName=self._fonts["oblique"],
            fontSize=9,
            leading=12,
            textColor=MUTED_COLOR,
            alignment=TA_CENTER,
            spaceBefore=4,
            spaceAfter=10,
        ))
        self._styles.add(ParagraphStyle(
            "FigureCaption",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=9,
            leading=12,
            textColor=TEXT_COLOR,
            alignment=TA_CENTER,
            spaceBefore=4,
            spaceAfter=12,
        ))
        self._styles.add(ParagraphStyle(
            "TocEntry1",
            parent=self._styles["Normal"],
            fontName=self._fonts["bold"],
            fontSize=11,
            leading=18,
            textColor=TEXT_COLOR,
            leftIndent=0,
        ))
        self._styles.add(ParagraphStyle(
            "TocEntry2",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=10,
            leading=16,
            textColor=MUTED_COLOR,
            leftIndent=20,
        ))
        self._styles.add(ParagraphStyle(
            "ConclusionText",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=10,
            leading=14,
            textColor=TEXT_COLOR,
            spaceAfter=8,
            leftIndent=10,
        ))
        self._styles.add(ParagraphStyle(
            "ReferenceText",
            parent=self._styles["Normal"],
            fontName=self._fonts["normal"],
            fontSize=9,
            leading=12,
            textColor=TEXT_COLOR,
            spaceAfter=4,
        ))

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def generate_report(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
        output_path: str,
    ) -> str:
        """Generate a PDF report and save to disk. Returns the output path."""
        pdf_bytes = self.generate_report_bytes(project_data, experiment_data)
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        with open(output_path, "wb") as fh:
            fh.write(pdf_bytes)
        return output_path

    def _reset(self):
        self._story = []
        self._toc_entries = []
        self._figures = []
        self._figure_counter = 0
        self._section_counter = 0

    def generate_report_bytes(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ) -> bytes:
        """Generate a PDF report and return as bytes."""
        self._reset()
        exp = self._canonicalize_experiment(experiment_data)

        title = project_data.get("name", "MatPilot Scientific Analysis Report")

        self._render_figures(exp)
        self._add_cover_page(project_data, exp)
        self._add_toc_placeholder()
        self._add_project_info(project_data)
        self._add_sample_info(exp)
        self._add_experimental_conditions(exp)
        self._add_data_summary(exp)
        self._add_processing_workflow(exp)
        self._add_phase_identification(exp)
        self._add_rietveld_summary(exp)
        self._add_refinement_statistics(exp)
        self._add_conclusions(project_data, exp)
        self._add_methodology(exp)
        self._add_references()
        self._add_figures_section()

        self._fill_toc()
        return self._build_pdf(title)

    def generate_docx_bytes(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ) -> bytes:
        """Generate a real .docx (OOXML) report and return as bytes."""
        if not _DOCX_AVAILABLE:
            raise RuntimeError(
                "python-docx is not installed. Run `pip install python-docx` and add it to requirements.txt."
            )
        self._reset()
        exp = self._canonicalize_experiment(experiment_data)
        self._render_figures(exp)

        doc = Document()
        normal = doc.styles["Normal"]
        normal.font.name = "Calibri"
        normal.font.size = _DocxPt(10.5)

        doc.add_heading("MatPilot Scientific Analysis Report", 0)
        p = doc.add_paragraph()
        run = p.add_run(str(project_data.get("name", "Untitled Project")))
        run.bold = True
        p = doc.add_paragraph()
        p.add_run("Generated by MatPilot \u2014 ").italic = True
        p.add_run(datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")).italic = True

        for section in self._collect_sections(exp, project_data):
            if section["kind"] == "heading":
                doc.add_heading(section["text"], level=section.get("level", 1))
            elif section["kind"] == "paragraph":
                doc.add_paragraph(section["text"])
            elif section["kind"] == "table":
                self._docx_add_table(
                    doc,
                    section["headers"],
                    section["rows"],
                    caption=section.get("caption"),
                )

        if self._figures:
            doc.add_heading("12. Figures", level=1)
            for fig in self._figures:
                p = doc.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run()
                run.add_picture(io.BytesIO(fig["data"]), width=_DocxInches(6.3))
                cap = doc.add_paragraph(f"Figure {fig['number']}: {fig['caption']}")
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if cap.runs:
                    cap.runs[0].italic = True

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.read()

    def generate_txt_bytes(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ) -> bytes:
        """Generate a plain-text report and return as bytes."""
        self._reset()
        exp = self._canonicalize_experiment(experiment_data)
        self._render_figures(exp)

        lines: List[str] = []
        lines.append("=" * 74)
        lines.append("MatPilot Scientific Analysis Report")
        lines.append("=" * 74)
        lines.append(f"Project: {project_data.get('name', 'Untitled Project')}")
        lines.append(f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}")
        lines.append("")

        for section in self._collect_sections(exp, project_data):
            if section["kind"] == "heading":
                lines.append("")
                lines.append(section["text"].upper())
                lines.append("-" * 74)
            elif section["kind"] == "table":
                lines.extend(self._txt_table(section["headers"], section["rows"]))
            elif section["kind"] == "paragraph":
                lines.append(section["text"])
                lines.append("")

        if self._figures:
            lines.append("")
            lines.append("APPENDIX: FIGURES")
            lines.append("-" * 74)
            for fig in self._figures:
                lines.append(f"Figure {fig['number']}: {fig['caption']} (rendered in PDF/DOCX/PPTX)")

        return ("\n".join(lines).rstrip() + "\n").encode("utf-8")

    def generate_pptx_bytes(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ) -> bytes:
        """Generate a PowerPoint report and return as bytes."""
        if not _PPTX_AVAILABLE:
            raise RuntimeError(
                "python-pptx is not installed. Run `pip install python-pptx` and add it to requirements.txt."
            )
        self._reset()
        exp = self._canonicalize_experiment(experiment_data)
        self._render_figures(exp)

        prs = Presentation()
        prs.slide_width = _PptxInches(13.333)
        prs.slide_height = _PptxInches(7.5)

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "MatPilot Scientific Analysis Report"
        subtitle = title_slide.placeholders[1]
        subtitle.text = (
            f"{project_data.get('name', 'Untitled Project')}\n"
            f"Generated by MatPilot \u2014 {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}"
        )

        current_title = ""
        for section in self._collect_sections(exp, project_data):
            if section["kind"] == "heading":
                current_title = section["text"]
                if section.get("level", 1) == 1:
                    s = prs.slides.add_slide(prs.slide_layouts[1])
                    s.shapes.title.text = current_title
                    self._pptx_text_slide = s
            elif section["kind"] == "paragraph":
                self._pptx_add_paragraph(section["text"])
            elif section["kind"] == "table":
                table_slide = prs.slides.add_slide(prs.slide_layouts[1])
                table_slide.shapes.title.text = current_title or "Results"
                self._pptx_add_table(
                    table_slide,
                    section["headers"],
                    section["rows"],
                    caption=section.get("caption"),
                )

        for fig in self._figures:
            fig_slide = prs.slides.add_slide(prs.slide_layouts[6])
            pic_width = _PptxInches(10.5)
            pic = fig_slide.shapes.add_picture(
                io.BytesIO(fig["data"]),
                _PptxInches(1.4),
                _PptxInches(0.6),
                width=pic_width,
            )
            tb = fig_slide.shapes.add_textbox(
                _PptxInches(1.4),
                _PptxInches(6.8),
                _PptxInches(10.5),
                _PptxInches(0.5),
            )
            tb.text_frame.text = f"Figure {fig['number']}: {fig['caption']}"
            tb.text_frame.paragraphs[0].font.size = _PptxPt(14)
            tb.text_frame.paragraphs[0].font.italic = True

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    # ------------------------------------------------------------------
    # Shared content model
    # ------------------------------------------------------------------

    def _collect_sections(
        self,
        exp: Dict[str, Any],
        project_data: Dict[str, Any],
    ) -> List[Dict[str, Any]]:
        """Return a flat, ordered list of heading/paragraph/table sections.

        Used by the DOCX, TXT and PPTX renderers. The PDF renderer mirrors the
        same section numbering.
        """
        sections: List[Dict[str, Any]] = []

        sections.append({"kind": "heading", "text": "1. Project Information", "level": 1})
        sections.append({
            "kind": "table",
            "headers": ["Property", "Value"],
            "rows": self._project_info_rows(project_data),
            "caption": "Table 1: Project information summary",
        })

        sections.append({"kind": "heading", "text": "2. Sample Information", "level": 1})
        sections.append({
            "kind": "table",
            "headers": ["Property", "Value"],
            "rows": self._sample_info_rows(exp),
            "caption": "Table 2: Sample and radiation information",
        })

        sections.append({"kind": "heading", "text": "3. Experimental Conditions", "level": 1})
        sections.append({
            "kind": "table",
            "headers": ["Parameter", "Value"],
            "rows": self._experimental_conditions_rows(exp),
            "caption": "Table 3: Experimental conditions",
        })

        sections.append({"kind": "heading", "text": "4. Data Summary", "level": 1})
        sections.append({
            "kind": "table",
            "headers": ["Metric", "Value"],
            "rows": self._data_summary_rows(exp),
            "caption": "Table 4: Data summary statistics",
        })

        peaks = exp.get("detected_peaks") or []
        if peaks:
            sections.append({"kind": "heading", "text": "4.1 Detected Peaks", "level": 2})
            sections.append({
                "kind": "table",
                "headers": self._peak_headers(),
                "rows": self._peak_rows(exp),
                "caption": f"Table 5: Top {min(20, len(peaks))} detected peaks",
            })

        sections.append({"kind": "heading", "text": "5. Processing Workflow", "level": 1})
        stages = exp.get("pipeline_stages") or []
        if stages:
            sections.append({
                "kind": "table",
                "headers": ["#", "Stage", "Status", "Duration (s)"],
                "rows": self._stage_rows(exp),
                "caption": "Table 6: Pipeline processing stages",
            })
        else:
            sections.append({
                "kind": "paragraph",
                "text": "No pipeline processing stages have been recorded for this experiment.",
            })

        sections.append({"kind": "heading", "text": "6. Phase Identification Results", "level": 1})
        phases = exp.get("candidate_phases") or []
        if phases:
            sections.append({
                "kind": "table",
                "headers": self._phase_headers(),
                "rows": self._phase_rows(exp),
                "caption": f"Table 7: Top {min(15, len(phases))} candidate phases",
            })
        else:
            sections.append({"kind": "paragraph", "text": "No phase identification results available."})

        sections.append({"kind": "heading", "text": "7. Rietveld Refinement Summary", "level": 1})
        rt = exp.get("rietveld_results")
        if rt:
            sections.append({
                "kind": "table",
                "headers": ["Refinement Metric", "Value"],
                "rows": self._rietveld_rows(exp),
                "caption": "Table 8: Rietveld refinement quality indicators",
            })
        else:
            sections.append({
                "kind": "paragraph",
                "text": "Rietveld refinement has not been performed for this experiment.",
            })

        sections.append({"kind": "heading", "text": "8. Refinement Statistics", "level": 1})
        if rt:
            phases_frac = [p for p in (rt.get("phases") or []) if p.get("fraction") is not None]
            refined = rt.get("refined_parameters") or []
            if phases_frac:
                sections.append({"kind": "heading", "text": "8.1 Phase Fractions", "level": 2})
                sections.append({
                    "kind": "table",
                    "headers": ["Phase", "Formula", "Fraction (%)"],
                    "rows": self._phase_fraction_rows(exp),
                    "caption": "Table 9: Phase fractions from Rietveld refinement",
                })
            if refined:
                sections.append({"kind": "heading", "text": "8.2 Refined Parameters", "level": 2})
                sections.append({
                    "kind": "table",
                    "headers": self._refined_param_headers(),
                    "rows": self._refined_param_rows(exp),
                    "caption": "Table 10: Refined structural parameters (with esds)",
                })
            if not phases_frac and not refined:
                sections.append({
                    "kind": "paragraph",
                    "text": "Detailed refinement statistics are not available.",
                })
        else:
            sections.append({"kind": "paragraph", "text": "No refinement statistics available."})

        sections.append({"kind": "heading", "text": "9. Scientific Conclusions", "level": 1})
        for conclusion in self._generate_conclusions(project_data, exp):
            sections.append({"kind": "paragraph", "text": conclusion})

        sections.append({"kind": "heading", "text": "10. Methodology", "level": 1})
        for paragraph in self._methodology_paragraphs(exp):
            sections.append({"kind": "paragraph", "text": paragraph})

        sections.append({"kind": "heading", "text": "11. References", "level": 1})
        for i, ref in enumerate(REFERENCES, 1):
            sections.append({"kind": "paragraph", "text": f"[{i}] {ref}"})

        return sections

    def _project_info_rows(self, project_data: Dict[str, Any]) -> List[List[str]]:
        return [
            ["Project Name", str(project_data.get("name", "N/A"))],
            ["Material", str(project_data.get("material", "N/A"))],
            ["Creation Date", str(project_data.get("created_at", "N/A"))],
            ["Status", str(project_data.get("status", "N/A"))],
        ]

    def _sample_info_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        wavelength = exp.get("wavelength", 1.5406)
        two_theta = exp.get("two_theta") or []

        two_theta_range = "N/A"
        if len(two_theta) >= 2:
            two_theta_range = f"{min(two_theta):.2f} \u2013 {max(two_theta):.2f}"

        return [
            ["Material Formula", str(exp.get("name") or "Unknown Material")],
            ["Radiation Source", str(exp.get("radiation_type_display") or "X-ray")],
            ["Wavelength (\u00c5)", _fmt(wavelength, ".4f")],
            ["2\u03b8 Range (\u00b0)", two_theta_range],
            ["Number of Data Points", str(len(two_theta)) if two_theta else "N/A"],
        ]

    def _experimental_conditions_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        two_theta = exp.get("two_theta") or []

        step_size = exp.get("step_size")
        if step_size is None and len(two_theta) >= 2:
            diffs = np.diff(two_theta)
            if len(diffs) and np.mean(diffs) > 0:
                step_size = float(np.mean(diffs))

        scan_range = exp.get("scan_range")
        scan_range_str = "N/A"
        if isinstance(scan_range, list) and len(scan_range) >= 2:
            scan_range_str = f"{scan_range[0]:.2f} \u2013 {scan_range[1]:.2f} \u00b0"

        instrument = str(exp.get("instrument") or "N/A")
        temperature = exp.get("temperature_k")
        temperature_str = "N/A (assumed ambient)" if temperature is None else f"{_fmt(temperature, '.1f')} K"

        scan_time = exp.get("scan_time_seconds")
        scan_time_str = "N/A" if scan_time is None else f"{_fmt(scan_time, '.2f')} s"

        rows = [
            ["Instrument", instrument],
            ["Radiation Type", str(exp.get("radiation_type_display") or "X-ray")],
            ["Data Points", str(len(two_theta)) if two_theta else "N/A"],
            ["Scan Range (2\u03b8)", scan_range_str],
            ["Step Size (\u00b0)", _fmt(step_size, ".4f")],
            ["Scan Time", scan_time_str],
            ["Temperature", temperature_str],
        ]
        notes = exp.get("notes")
        if notes:
            rows.append(["Notes", str(notes)])
        return rows

    def _data_summary_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        intensity = exp.get("intensity") or []
        detected_peaks = exp.get("detected_peaks") or []

        snr = "N/A"
        bg_mean = "N/A"
        bg_std = "N/A"
        if len(intensity) > 0:
            arr = np.array(intensity, dtype=float)
            bg_est = float(np.percentile(arr, 10))
            signal_est = float(np.percentile(arr, 95))
            noise_est = float(np.std(arr[:max(1, len(arr) // 20)]))
            if noise_est > 0:
                snr = f"{signal_est / noise_est:.1f}"
            bg_mean = f"{bg_est:.2f}"
            bg_std = f"{noise_est:.2f}"

        return [
            ["Number of Peaks Detected", str(len(detected_peaks))],
            ["Signal-to-Noise Ratio", snr],
            ["Background Mean (counts)", bg_mean],
            ["Background Std Dev", bg_std],
            ["Total Data Points", str(len(intensity)) if intensity else "N/A"],
        ]

    def _peak_headers(self) -> List[str]:
        return ["#", "2\u03b8 (\u00b0)", "d-spacing (\u00c5)", "Rel. I (%)", "Intensity", "FWHM (\u00b0)", "Unc. (\u00b0)"]

    def _peak_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        peaks = exp.get("detected_peaks") or []
        rows = []
        for idx, peak in enumerate(peaks[:20], 1):
            rows.append([
                str(idx),
                _fmt(peak.get("two_theta"), ".3f"),
                _fmt(peak.get("d_spacing"), ".4f"),
                _fmt(peak.get("relative_intensity"), ".1f"),
                _fmt(peak.get("intensity"), ".1f"),
                _fmt(peak.get("fwhm"), ".4f"),
                _fmt(peak.get("esd"), ".4f"),
            ])
        return rows

    def _stage_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        stages = exp.get("pipeline_stages") or []
        rows = []
        for idx, stage in enumerate(stages, 1):
            stage_name = stage.get("name", "Unknown")
            display_name = PIPELINE_STAGE_LABELS.get(stage_name, stage_name)
            status = stage.get("status", "pending")
            duration = stage.get("duration_seconds")
            duration_str = _fmt(duration, ".2f") if duration is not None else "N/A"
            status_display = status.capitalize() if isinstance(status, str) else str(status)
            rows.append([str(idx), display_name, status_display, duration_str])
        return rows

    def _phase_headers(self) -> List[str]:
        return ["Rank", "Material Name", "Formula", "Match Score", "FOM", "Matched Peaks", "Confidence"]

    def _phase_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        phases = exp.get("candidate_phases") or []
        rows = []
        for phase in phases[:15]:
            matched = phase.get("matched_peaks")
            total = phase.get("total_reference_peaks")
            if matched is not None and total is not None:
                matched_str = f"{_fmt_int(matched)}/{_fmt_int(total)}"
            elif matched is not None:
                matched_str = _fmt_int(matched)
            else:
                matched_str = "N/A"
            rows.append([
                str(phase.get("rank", "\u2014")),
                str(phase.get("material_name", "Unknown")),
                str(phase.get("material_formula", "\u2014")),
                _fmt_pct(phase.get("match_score")),
                _fmt(phase.get("fom"), ".3f"),
                matched_str,
                str(phase.get("confidence", "\u2014")),
            ])
        return rows

    def _rietveld_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        rt = exp.get("rietveld_results")
        if not rt:
            return []
        return [
            ["R_wp (%)", _fmt(rt.get("r_wp"), ".4f")],
            ["R_p (%)", _fmt(rt.get("r_p"), ".4f")],
            ["R_exp (%)", _fmt(rt.get("r_exp"), ".4f")],
            ["Chi-squared (\u03c7\u00b2)", _fmt(rt.get("chi_squared"), ".4f")],
            ["Goodness of Fit (GoF)", _fmt(rt.get("gof"), ".4f")],
            ["Iterations", _fmt_int(rt.get("iterations"))],
            ["Phases Used", _fmt_int(rt.get("phases_used_count"))],
            ["Status", str(rt.get("status") or "N/A")],
        ]

    def _phase_fraction_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        rt = exp.get("rietveld_results")
        if not rt:
            return []
        rows = []
        for phase in rt.get("phases") or []:
            if phase.get("fraction") is None:
                continue
            rows.append([
                str(phase.get("name") or "Unknown"),
                str(phase.get("formula") or "\u2014"),
                _fmt_pct(phase.get("fraction")),
            ])
        return rows

    def _refined_param_headers(self) -> List[str]:
        return ["Parameter", "Refined Value", "Estimated \u03c3"]

    def _refined_param_rows(self, exp: Dict[str, Any]) -> List[List[str]]:
        rt = exp.get("rietveld_results")
        if not rt:
            return []
        rows = []
        for param in rt.get("refined_parameters") or []:
            rows.append([
                str(param.get("name") or "\u2014"),
                _fmt(param.get("refined"), ".6f"),
                _fmt(param.get("uncertainty"), ".6f"),
            ])
        return rows

    # ------------------------------------------------------------------
    # Section builders (PDF)
    # ------------------------------------------------------------------

    def _add_cover_page(self, project_data: Dict[str, Any], experiment_data: Dict[str, Any]):
        cover_elements: List[Any] = []
        cover_elements.append(Spacer(1, 4 * cm))
        cover_elements.append(Paragraph(
            "MatPilot Scientific<br/>Analysis Report",
            self._styles["CoverTitle"],
        ))
        cover_elements.append(Spacer(1, 0.8 * cm))

        project_name = project_data.get("name", "Untitled Project")
        cover_elements.append(Paragraph(project_name, self._styles["CoverSubtitle"]))
        cover_elements.append(Spacer(1, 0.5 * cm))

        material = project_data.get("material", "")
        if material:
            cover_elements.append(Paragraph(
                f"Material: {material}",
                self._styles["CoverInfo"],
            ))
            cover_elements.append(Spacer(1, 0.3 * cm))

        date_str = project_data.get("created_at", "")
        if date_str:
            cover_elements.append(Paragraph(
                f"Created: {date_str}",
                self._styles["CoverInfo"],
            ))
        else:
            cover_elements.append(Paragraph(
                f"Generated: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}",
                self._styles["CoverInfo"],
            ))

        cover_elements.append(Spacer(1, 0.3 * cm))
        status = project_data.get("status", "")
        if status:
            cover_elements.append(Paragraph(
                f"Status: {status}",
                self._styles["CoverInfo"],
            ))

        cover_elements.append(Spacer(1, 2 * cm))
        cover_elements.append(Paragraph(
            "Automated XRD Analysis Report",
            self._styles["CoverInfo"],
        ))
        cover_elements.append(Paragraph(
            "Generated by MatPilot",
            self._styles["CoverInfo"],
        ))

        self._story.extend(cover_elements)
        self._story.append(NextPageTemplate("content"))
        self._story.append(PageBreak())

    def _add_toc_placeholder(self):
        self._toc_heading_index = len(self._story)
        self._story.append(Paragraph(
            "Table of Contents",
            self._styles["SectionHeading"],
        ))
        self._story.append(Spacer(1, 0.5 * cm))
        self._toc_placeholder_index = len(self._story)
        self._story.append(Spacer(1, 0.1 * cm))
        self._story.append(PageBreak())

    def _fill_toc(self):
        toc_items: List[Any] = []
        for entry in self._toc_entries:
            style_name = "TocEntry1" if entry["level"] == 1 else "TocEntry2"
            page_label = str(entry.get("page", ""))
            dots = " " + "." * max(2, 60 - len(entry["text"]) - len(page_label))
            text = f"{entry['text']}{dots}{page_label}"
            toc_items.append(Paragraph(text, self._styles[style_name]))
        if not toc_items:
            toc_items.append(Paragraph("\u2014", self._styles["BodyText2"]))
        self._story[self._toc_placeholder_index:self._toc_placeholder_index + 1] = toc_items

    def _record_toc(self, text: str, level: int = 1):
        self._toc_entries.append({
            "text": text,
            "level": level,
            "page": "\u2014",
        })

    def _next_section_number(self) -> str:
        self._section_counter += 1
        return str(self._section_counter)

    def _add_section_heading(self, text: str, level: int = 1) -> str:
        section_num = self._next_section_number()
        heading_text = f"{section_num}. {text}"
        self._story.append(Spacer(1, 0.3 * cm))
        self._story.append(Paragraph(heading_text, self._styles["SectionHeading"]))
        self._record_toc(heading_text, level=level)
        self._story.append(Spacer(1, 0.2 * cm))
        return section_num

    def _add_subsection_heading(self, section_num: str, text: str):
        sub_text = f"{section_num}.{text}"
        self._story.append(Paragraph(sub_text, self._styles["SubHeading"]))
        self._record_toc(sub_text, level=2)

    def _add_body(self, text: str):
        self._story.append(Paragraph(text, self._styles["BodyText2"]))

    def _add_paragraph(self, text: str, style: str = "BodyText2"):
        self._story.append(Paragraph(text, self._styles[style]))

    def _build_table(
        self,
        headers: List[str],
        rows: List[List[str]],
        col_widths: Optional[List[float]] = None,
        caption: Optional[str] = None,
    ):
        data = [headers] + rows
        table = Table(data, colWidths=col_widths, repeatRows=1)

        style_cmds = [
            ("BACKGROUND", (0, 0), (-1, 0), TABLE_HEADER_BG),
            ("TEXTCOLOR", (0, 0), (-1, 0), TABLE_HEADER_TEXT),
            ("FONTNAME", (0, 0), (-1, 0), self._fonts["bold"]),
            ("FONTSIZE", (0, 0), (-1, 0), 9),
            ("FONTNAME", (0, 1), (-1, -1), self._fonts["normal"]),
            ("FONTSIZE", (0, 1), (-1, -1), 9),
            ("LEADING", (0, 0), (-1, -1), 13),
            ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("GRID", (0, 0), (-1, -1), 0.5, BORDER_COLOR),
            ("LINEBELOW", (0, 0), (-1, 0), 1.5, ACCENT_COLOR),
        ]

        for row_idx in range(1, len(data)):
            if row_idx % 2 == 0:
                style_cmds.append(
                    ("BACKGROUND", (0, row_idx), (-1, row_idx), TABLE_ALT_ROW)
                )

        table.setStyle(TableStyle(style_cmds))
        self._story.append(table)

        if caption:
            self._story.append(Paragraph(caption, self._styles["TableCaption"]))

    # ------------------------------------------------------------------
    # Figures
    # ------------------------------------------------------------------

    def _render_figures(self, exp: Dict[str, Any]):
        """Render all report figures into PNG bytes (matplotlib, Agg backend)."""
        two_theta = exp.get("two_theta") or []
        intensity = exp.get("intensity") or []

        if two_theta and intensity:
            self._create_figure(
                two_theta,
                intensity,
                "Experimental Diffraction Pattern",
                caption="Experimental XRD diffraction pattern",
            )

        processed_2t = exp.get("processed_two_theta") or []
        processed_int = exp.get("processed_intensity") or []
        if processed_2t and processed_int:
            identical = (
                list(processed_2t) == list(two_theta)
                and list(processed_int) == list(intensity)
            )
            if not identical:
                self._create_figure(
                    processed_2t,
                    processed_int,
                    "Background-Corrected Pattern",
                    caption="Background-corrected diffraction pattern",
                )

        rietveld = exp.get("rietveld_results")
        if rietveld:
            calc_2t = rietveld.get("calculated_two_theta") or []
            calc_int = rietveld.get("calculated_intensity") or []
            diff_int = rietveld.get("difference_intensity") or []

            if calc_2t and calc_int and two_theta:
                x_ref = two_theta if len(two_theta) == len(calc_2t) else two_theta[:len(calc_2t)]
                y_exp = intensity[:len(x_ref)] if len(intensity) >= len(x_ref) else intensity

                extra = []
                if len(calc_int) >= len(x_ref):
                    extra.append({
                        "x": x_ref,
                        "y": calc_int[:len(x_ref)],
                        "label": "Calculated",
                    })
                if diff_int and len(diff_int) >= len(x_ref):
                    extra.append({
                        "x": x_ref,
                        "y": diff_int[:len(x_ref)],
                        "label": "Difference",
                    })

                self._create_figure(
                    x_ref,
                    y_exp[:len(x_ref)],
                    "Rietveld Refinement: Experimental vs Calculated",
                    extra_traces=extra,
                    caption="Overlay of experimental, calculated, and difference patterns from Rietveld refinement",
                )

    def _create_figure(
        self,
        x_data: List[float],
        y_data: List[float],
        title: str,
        xlabel: str = r"2$\theta$ (degrees)",
        ylabel: str = "Intensity (a.u.)",
        extra_traces: Optional[List[Dict[str, Any]]] = None,
        caption: Optional[str] = None,
    ) -> bytes:
        fig, ax = plt.subplots(figsize=(FIGURE_WIDTH_CM / 2.54, FIGURE_HEIGHT_CM / 2.54), dpi=150)

        ax.plot(
            x_data, y_data,
            color="#f97316",
            linewidth=0.8,
            label="Experimental",
        )

        if extra_traces:
            trace_colors = ["#3b82f6", "#10b981", "#ef4444", "#8b5cf6"]
            for idx, trace in enumerate(extra_traces):
                color = trace_colors[idx % len(trace_colors)]
                ax.plot(
                    trace.get("x", []),
                    trace.get("y", []),
                    color=color,
                    linewidth=0.8,
                    label=trace.get("label", ""),
                )

        ax.set_xlabel(xlabel, fontsize=10, fontweight="bold")
        ax.set_ylabel(ylabel, fontsize=10, fontweight="bold")
        ax.set_title(title, fontsize=11, fontweight="bold", pad=8)
        ax.legend(fontsize=8, loc="upper right", framealpha=0.9)
        ax.grid(True, alpha=0.3, linestyle="--")
        ax.tick_params(labelsize=9)
        ax.xaxis.set_major_locator(ticker.MaxNLocator(8))
        ax.yaxis.set_major_locator(ticker.MaxNLocator(6))

        for spine in ax.spines.values():
            spine.set_linewidth(0.8)

        fig.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
        buf.seek(0)

        self._figure_counter += 1
        self._figures.append({
            "number": self._figure_counter,
            "data": buf.read(),
            "caption": caption or title,
        })
        return b""

    def _embed_figures(self):
        for fig_info in self._figures:
            buf = io.BytesIO(fig_info["data"])
            img = Image(buf, width=FIGURE_WIDTH_CM * cm, height=FIGURE_HEIGHT_CM * cm)
            self._story.append(Spacer(1, 0.3 * cm))
            self._story.append(img)
            caption_text = f"Figure {fig_info['number']}: {fig_info['caption']}"
            self._story.append(Paragraph(caption_text, self._styles["FigureCaption"]))
            self._story.append(Spacer(1, 0.3 * cm))

    # ------------------------------------------------------------------
    # Section 1: Project Information
    # ------------------------------------------------------------------

    def _add_project_info(self, experiment_data: Dict[str, Any]):
        self._add_section_heading("Project Information")
        self._add_property_table(
            self._project_info_rows(experiment_data),
            caption="Table 1: Project information summary",
        )

    def _add_property_table(self, rows: List[List[str]], caption: str):
        usable_width = self._page_width - self._margin_left - self._margin_right
        col_widths = [usable_width * 0.35, usable_width * 0.65]
        self._build_table(
            headers=["Property", "Value"],
            rows=rows,
            col_widths=col_widths,
            caption=caption,
        )

    # ------------------------------------------------------------------
    # Section 2: Sample Information
    # ------------------------------------------------------------------

    def _add_sample_info(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Sample Information")
        self._add_property_table(
            self._sample_info_rows(experiment_data),
            caption="Table 2: Sample and radiation information",
        )

    # ------------------------------------------------------------------
    # Section 3: Experimental Conditions
    # ------------------------------------------------------------------

    def _add_experimental_conditions(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Experimental Conditions")
        usable_width = self._page_width - self._margin_left - self._margin_right
        col_widths = [usable_width * 0.35, usable_width * 0.65]
        self._build_table(
            headers=["Parameter", "Value"],
            rows=self._experimental_conditions_rows(experiment_data),
            col_widths=col_widths,
            caption="Table 3: Experimental conditions",
        )

    # ------------------------------------------------------------------
    # Section 4: Data Summary
    # ------------------------------------------------------------------

    def _add_data_summary(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Data Summary")
        usable_width = self._page_width - self._margin_left - self._margin_right
        col_widths = [usable_width * 0.45, usable_width * 0.55]
        self._build_table(
            headers=["Metric", "Value"],
            rows=self._data_summary_rows(experiment_data),
            col_widths=col_widths,
            caption="Table 4: Data summary statistics",
        )

        peaks = experiment_data.get("detected_peaks") or []
        if peaks:
            self._add_body("")
            self._add_subsection_heading(section_num, "Detected Peaks")

            peak_rows = self._peak_rows(experiment_data)
            peak_col_widths = [
                usable_width * 0.06,
                usable_width * 0.16,
                usable_width * 0.15,
                usable_width * 0.13,
                usable_width * 0.14,
                usable_width * 0.16,
                usable_width * 0.20,
            ]
            self._build_table(
                headers=self._peak_headers(),
                rows=peak_rows,
                col_widths=peak_col_widths,
                caption=f"Table 5: Top {min(20, len(peaks))} detected peaks",
            )

    # ------------------------------------------------------------------
    # Section 5: Processing Workflow
    # ------------------------------------------------------------------

    def _add_processing_workflow(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Processing Workflow")

        pipeline_stages = experiment_data.get("pipeline_stages", [])
        if not pipeline_stages:
            self._add_body("No pipeline processing stages have been recorded for this experiment.")
            return

        stage_rows = self._stage_rows(experiment_data)
        total_duration = sum(
            s.get("duration_seconds", 0)
            for s in pipeline_stages
            if s.get("duration_seconds") is not None
        )

        usable_width = self._page_width - self._margin_left - self._margin_right
        stage_col_widths = [
            usable_width * 0.08,
            usable_width * 0.42,
            usable_width * 0.25,
            usable_width * 0.25,
        ]
        self._build_table(
            headers=["#", "Stage", "Status", "Duration (s)"],
            rows=stage_rows,
            col_widths=stage_col_widths,
            caption=f"Table 6: Pipeline processing stages (total: {_fmt(total_duration, '.2f')} s)",
        )

    # ------------------------------------------------------------------
    # Section 6: Phase Identification Results
    # ------------------------------------------------------------------

    def _add_phase_identification(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Phase Identification Results")

        candidate_phases = experiment_data.get("candidate_phases", [])
        if not candidate_phases:
            self._add_body("No phase identification results available.")
            return

        usable_width = self._page_width - self._margin_left - self._margin_right
        phase_col_widths = [
            usable_width * 0.06,
            usable_width * 0.24,
            usable_width * 0.20,
            usable_width * 0.14,
            usable_width * 0.12,
            usable_width * 0.14,
            usable_width * 0.10,
        ]
        self._build_table(
            headers=self._phase_headers(),
            rows=self._phase_rows(experiment_data),
            col_widths=phase_col_widths,
            caption=f"Table 7: Top {min(15, len(candidate_phases))} candidate phases",
        )

    # ------------------------------------------------------------------
    # Section 7: Rietveld Refinement Summary
    # ------------------------------------------------------------------

    def _add_rietveld_summary(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Rietveld Refinement Summary")

        rietveld = experiment_data.get("rietveld_results")
        if not rietveld:
            self._add_body("Rietveld refinement has not been performed for this experiment.")
            return

        usable_width = self._page_width - self._margin_left - self._margin_right
        col_widths = [usable_width * 0.40, usable_width * 0.60]
        self._build_table(
            headers=["Refinement Metric", "Value"],
            rows=self._rietveld_rows(experiment_data),
            col_widths=col_widths,
            caption="Table 8: Rietveld refinement quality indicators",
        )

    # ------------------------------------------------------------------
    # Section 8: Refinement Statistics
    # ------------------------------------------------------------------

    def _add_refinement_statistics(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Refinement Statistics")

        rietveld = experiment_data.get("rietveld_results")
        if not rietveld:
            self._add_body("No refinement statistics available.")
            return

        phases_frac = [p for p in (rietveld.get("phases") or []) if p.get("fraction") is not None]
        refined_params = rietveld.get("refined_parameters") or []

        if phases_frac:
            self._add_subsection_heading(section_num, "Phase Fractions")
            usable_width = self._page_width - self._margin_left - self._margin_right
            pf_col_widths = [
                usable_width * 0.35,
                usable_width * 0.35,
                usable_width * 0.30,
            ]
            self._build_table(
                headers=["Phase", "Formula", "Fraction (%)"],
                rows=self._phase_fraction_rows(experiment_data),
                col_widths=pf_col_widths,
                caption="Table 9: Phase fractions from Rietveld refinement",
            )

        if refined_params:
            self._add_subsection_heading(section_num, "Refined Parameters")
            usable_width = self._page_width - self._margin_left - self._margin_right
            rp_col_widths = [usable_width * 0.50, usable_width * 0.25, usable_width * 0.25]
            self._build_table(
                headers=self._refined_param_headers(),
                rows=self._refined_param_rows(experiment_data),
                col_widths=rp_col_widths,
                caption="Table 10: Refined structural parameters (with esds)",
            )

        if not phases_frac and not refined_params:
            self._add_body("Detailed refinement statistics are not available.")

    # ------------------------------------------------------------------
    # Section 9: Conclusions
    # ------------------------------------------------------------------

    def _add_conclusions(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ):
        section_num = self._add_section_heading("Scientific Conclusions")

        conclusions = self._generate_conclusions(project_data, experiment_data)
        for conclusion in conclusions:
            self._add_paragraph(conclusion, "ConclusionText")

    def _generate_conclusions(
        self,
        project_data: Dict[str, Any],
        experiment_data: Dict[str, Any],
    ) -> List[str]:
        conclusions: List[str] = []
        material = project_data.get("material", "the sample")

        two_theta = experiment_data.get("two_theta") or []
        intensity = experiment_data.get("intensity") or []
        detected_peaks = experiment_data.get("detected_peaks") or []
        candidate_phases = experiment_data.get("candidate_phases") or []
        rietveld = experiment_data.get("rietveld_results")

        if two_theta and intensity:
            conclusions.append(
                f"The diffraction data for <b>{material}</b> consists of "
                f"<b>{len(two_theta)}</b> data points spanning the 2\u03b8 range of "
                f"<b>{min(two_theta):.2f}\u00b0 \u2013 {max(two_theta):.2f}\u00b0</b>."
            )

        if detected_peaks:
            conclusions.append(
                f"A total of <b>{len(detected_peaks)}</b> diffraction peaks were "
                f"identified in the experimental pattern."
            )

        if candidate_phases:
            top = candidate_phases[0]
            name = top.get("material_name", "Unknown")
            score = top.get("match_score")
            score_str = f"{float(score):.1%}" if score is not None else "N/A"
            conclusions.append(
                f"The best-matching phase from the database search is "
                f"<b>{name}</b> with a match score of <b>{score_str}</b>."
            )
            if len(candidate_phases) > 1:
                runner_up = candidate_phases[1]
                runner_name = runner_up.get("material_name", "Unknown")
                conclusions.append(
                    f"Additional candidate phases include <b>{runner_name}</b> "
                    f"and {len(candidate_phases) - 1} other matches."
                )

        if rietveld:
            gof = rietveld.get("gof")
            chi2 = rietveld.get("chi_squared")
            rwp = rietveld.get("r_wp")
            if gof is not None:
                gof_val = float(gof)
                if gof_val < 1.5:
                    quality = "good"
                elif gof_val < 2.5:
                    quality = "acceptable"
                else:
                    quality = "poor"
                conclusions.append(
                    f"The Rietveld refinement yielded a goodness-of-fit (GoF) of "
                    f"<b>{gof_val:.4f}</b>, indicating a <b>{quality}</b> fit quality."
                )
            if rwp is not None:
                conclusions.append(
                    f"The weighted profile R-factor (R_wp) is "
                    f"<b>{float(rwp):.4f}%</b>."
                )

            phases_used = rietveld.get("phases_used_count")
            phases = rietveld.get("phases") or []
            if phases:
                phase_names = [p.get("name", "Unknown") for p in phases[:3]]
                conclusions.append(
                    f"The refinement utilized <b>{phases_used if phases_used is not None else len(phases)}</b> phase(s): "
                    f"<b>{', '.join(phase_names)}</b>."
                )

        if not conclusions:
            conclusions.append(
                "Insufficient data is available to draw scientific conclusions. "
                "Please complete the analysis pipeline and Rietveld refinement."
            )

        conclusions.append(
            "This report was automatically generated by the MatPilot scientific "
            "analysis platform. All results should be reviewed and validated by "
            "the researcher before publication."
        )

        return conclusions

    # ------------------------------------------------------------------
    # Section 10: Methodology
    # ------------------------------------------------------------------

    def _add_methodology(self, experiment_data: Dict[str, Any]):
        section_num = self._add_section_heading("Methodology")
        for paragraph in self._methodology_paragraphs(experiment_data):
            self._add_paragraph(paragraph, "BodyText2")

    def _methodology_paragraphs(self, exp: Dict[str, Any]) -> List[str]:
        paragraphs: List[str] = []

        instrument = str(exp.get("instrument") or "").strip()
        instrument_display = instrument if instrument not in ("", "N/A") else "a laboratory diffractometer"
        radiation = str(exp.get("radiation_type_display") or "X-ray")
        wl = _safe_float(exp.get("wavelength"), None)
        wl_s = _fmt(wl, ".4f")

        paragraphs.append(
            "Powder X-ray diffraction (XRD) data were collected using the MatPilot "
            f"analysis platform on {instrument_display} with {radiation} radiation "
            f"(\u03bb = {wl_s} \u00c5)."
        )

        scan_range = exp.get("scan_range")
        step = exp.get("step_size")
        if isinstance(scan_range, list) and len(scan_range) >= 2:
            range_s = f"{scan_range[0]:.2f}\u00b0 to {scan_range[1]:.2f}\u00b0"
            if step is not None:
                paragraphs.append(
                    f"Diffraction data were recorded over the 2\u03b8 range of {range_s} "
                    f"with a step size of {float(step):.4f}\u00b0."
                )
            else:
                paragraphs.append(f"Diffraction data were recorded over the 2\u03b8 range of {range_s}.")

        stages = exp.get("pipeline_stages") or []
        if stages:
            names = [PIPELINE_STAGE_LABELS.get(s.get("name", ""), s.get("name", "")) for s in stages]
            paragraphs.append(
                "The raw diffraction pattern was processed through the following pipeline "
                f"stages, in sequence: {'; '.join(names)}."
            )
        else:
            paragraphs.append(
                "The raw diffraction pattern was used without additional processing."
            )

        peaks = exp.get("detected_peaks") or []
        if peaks:
            paragraphs.append(
                f"Peak detection identified {len(peaks)} diffraction maxima. Peak positions, "
                "d-spacings, relative intensities, full widths at half maximum (FWHM), and "
                "estimated positional uncertainties are tabulated in Section 4.1."
            )

        phases = exp.get("candidate_phases") or []
        if phases:
            paragraphs.append(
                "Phase identification was performed by matching the experimental peak positions "
                "and intensities against reference patterns from the Crystallography Open Database "
                "(COD) and the local MatPilot reference database. Candidates were ranked by a "
                "combined match score incorporating peak-position quality (figure of merit), "
                "peak fraction, and intensity (cosine) similarity."
            )

        rietveld = exp.get("rietveld_results")
        if rietveld:
            paragraphs.append(
                "Rietveld refinement was performed using a pseudo-Voigt peak profile model and a "
                "polynomial background. Refined parameters included the scale factor, zero shift, "
                "peak-shape parameters, and phase fractions. The quality of the fit is reported via "
                "the weighted profile R-factor (R_wp), the expected R-factor (R_exp), the "
                "chi-squared statistic, and the goodness-of-fit (GoF)."
            )

        if not stages and not phases and not rietveld:
            paragraphs.append(
                "This experiment has not yet completed the full analysis pipeline; only the raw "
                "diffraction data are reported here."
            )

        return paragraphs

    # ------------------------------------------------------------------
    # Section 11: References
    # ------------------------------------------------------------------

    def _add_references(self):
        section_num = self._add_section_heading("References")
        for i, ref in enumerate(REFERENCES, 1):
            self._add_paragraph(f"[{i}] {ref}", "ReferenceText")

    # ------------------------------------------------------------------
    # Section 12: Figures (appendix)
    # ------------------------------------------------------------------

    def _add_figures_section(self):
        if not self._figures:
            return
        self._story.append(PageBreak())
        appendix_num = self._next_section_number()
        heading_text = f"{appendix_num}. Figures"
        self._story.append(Paragraph(heading_text, self._styles["SectionHeading"]))
        self._record_toc(heading_text, level=1)
        self._story.append(Spacer(1, 0.2 * cm))
        self._embed_figures()

    # ------------------------------------------------------------------
    # DOCX helpers
    # ------------------------------------------------------------------

    def _docx_add_table(
        self,
        doc,
        headers: List[str],
        rows: List[List[str]],
        caption: Optional[str] = None,
    ):
        if not rows:
            return
        table = doc.add_table(rows=1 + len(rows), cols=len(headers))
        table.style = "Table Grid"
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr_cells[i].text = str(h)
            for paragraph in hdr_cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        for r_idx, row in enumerate(rows, start=1):
            cells = table.rows[r_idx].cells
            for c_idx in range(len(headers)):
                value = row[c_idx] if c_idx < len(row) else ""
                cells[c_idx].text = str(value)
        if caption:
            cap = doc.add_paragraph(str(caption))
            cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if cap.runs:
                cap.runs[0].italic = True

    # ------------------------------------------------------------------
    # TXT helpers
    # ------------------------------------------------------------------

    def _txt_table(self, headers: List[str], rows: List[List[str]]) -> List[str]:
        data = [headers] + rows
        ncols = len(headers)
        widths = []
        for c in range(ncols):
            col = [str(r[c]) if c < len(r) else "" for r in data]
            widths.append(max(len(x) for x in col))

        def fmt_row(row: List[str]) -> str:
            return "  ".join(
                str(row[c]).ljust(widths[c]) if c < len(row) else "".ljust(widths[c])
                for c in range(ncols)
            ).rstrip()

        lines = [fmt_row(headers), "  ".join("-" * w for w in widths)]
        for row in rows:
            lines.append(fmt_row(row))
        lines.append("")
        return lines

    # ------------------------------------------------------------------
    # PPTX helpers
    # ------------------------------------------------------------------

    def _pptx_add_paragraph(self, text: str):
        slide = getattr(self, "_pptx_text_slide", None)
        if slide is None:
            return
        try:
            body = slide.placeholders[1]
            tf = body.text_frame
            if tf.paragraphs and len(tf.paragraphs) == 1 and not tf.text.strip():
                tf.text = text
                tf.paragraphs[0].font.size = _PptxPt(14)
                return
            p = tf.add_paragraph()
            p.text = text
            p.font.size = _PptxPt(14)
        except Exception:
            pass

    def _pptx_add_table(self, slide, headers, rows, caption=None):
        data = [headers] + rows
        n_rows = len(data)
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            _PptxInches(0.5), _PptxInches(1.2),
            _PptxInches(12.3), _PptxInches(0.3 * n_rows + 0.5),
        )
        table = table_shape.table
        for r in range(n_rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                value = data[r][c] if c < len(data[r]) else ""
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = _PptxPt(11)
                if r == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
        if caption:
            tb = slide.shapes.add_textbox(
                _PptxInches(0.5), _PptxInches(6.9), _PptxInches(12.3), _PptxInches(0.4)
            )
            tb.text_frame.text = str(caption)
            tb.text_frame.paragraphs[0].font.size = _PptxPt(10)
            tb.text_frame.paragraphs[0].font.italic = True

    # ------------------------------------------------------------------
    # PDF assembly
    # ------------------------------------------------------------------

    def _build_pdf(self, title: str) -> bytes:
        buf = io.BytesIO()

        doc = BaseDocTemplate(
            buf,
            pagesize=A4,
            leftMargin=self._margin_left,
            rightMargin=self._margin_right,
            topMargin=self._margin_top,
            bottomMargin=self._margin_bottom,
            title=f"MatPilot Report \u2013 {title}",
            author="MatPilot",
            pageCompression=0,
        )

        frame_cover = Frame(
            self._margin_left,
            self._margin_bottom,
            self._page_width - self._margin_left - self._margin_right,
            self._page_height - self._margin_top - self._margin_bottom,
            id="cover",
        )
        frame_content = Frame(
            self._margin_left,
            self._margin_bottom + 0.5 * cm,
            self._page_width - self._margin_left - self._margin_right,
            self._page_height - self._margin_top - self._margin_bottom - 0.5 * cm,
            id="content",
        )

        def _cover_bg(canvas, doc):
            canvas.saveState()
            canvas.setFillColor(HEADER_BG)
            canvas.rect(0, 0, self._page_width, self._page_height, fill=True, stroke=False)

            canvas.setFillColor(colors.HexColor("#f97316"))
            canvas.rect(
                0,
                self._page_height - 0.4 * cm,
                self._page_width,
                0.4 * cm,
                fill=True,
                stroke=False,
            )
            canvas.rect(
                0,
                0,
                self._page_width,
                0.4 * cm,
                fill=True,
                stroke=False,
            )
            canvas.restoreState()

        doc.addPageTemplates([
            PageTemplate(id="cover", frames=[frame_cover], onPage=_cover_bg),
            PageTemplate(
                id="content",
                frames=[frame_content],
                onPage=_NumberedCanvas(doc, title, self._fonts),
            ),
        ])

        doc.build(self._story)
        buf.seek(0)
        return buf.read()


# ----------------------------------------------------------------------
# Module-level convenience API
# ----------------------------------------------------------------------


def generate_report_bytes(
    report_format: str,
    project_data: Dict[str, Any],
    experiment_data: Dict[str, Any],
) -> bytes:
    """Generate a report in the requested format. Returns raw bytes."""
    generator = ReportGenerator()
    fmt = str(report_format).lower()
    if fmt == "pdf":
        return generator.generate_report_bytes(project_data, experiment_data)
    if fmt == "docx":
        return generator.generate_docx_bytes(project_data, experiment_data)
    if fmt == "txt":
        return generator.generate_txt_bytes(project_data, experiment_data)
    if fmt == "pptx":
        return generator.generate_pptx_bytes(project_data, experiment_data)
    raise ValueError(f"Unsupported report format: {report_format}")
